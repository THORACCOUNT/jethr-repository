from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

FONT = "Courier New"


def customize_text(text_frame, font_name, font_size, color, bold=False, underline=False, spacing=None):
    """Method used to customize text with font, color, space and other attributes"""
    for paragraph in text_frame.paragraphs:
        if spacing:
            paragraph.space_after = spacing
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.underline = underline


class JetHRPresentation:
    """Class for customize a presentation object from pptx"""

    def __init__(self, output_path):
        self.prs = Presentation()
        # Set layout for slides
        self.title_layout = self.prs.slide_layouts[0]
        self.content_layout = self.prs.slide_layouts[1]
        # Set path folder for presentation output
        self.output_path = output_path
        # Set font for slides
        self.title_font = (FONT, Pt(44))
        self.content_font = (FONT, Pt(22))
        # Set colors for slides and their content
        self.bg_color = RGBColor(50, 50, 50)
        self.title_color = RGBColor(0, 0, 0)
        self.content_color = RGBColor(224, 219, 209)

    def add_background(self, slide, color):
        """Method to add background to slides"""
        # Set initial position of the rectangles
        left = top = Inches(0)
        # Get dimensions of current slide
        width = self.prs.slide_width
        height = self.prs.slide_height
        # Add a rectangle that covers all slide
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        # Remove rectangle edge
        background.line.fill.background()
        # Set rectangle filling
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
        # Move rectangle behind other elements, since by default it is superimposed
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)

    def add_title_slide(self, title, subtitle):
        """Method that add custom title to main slide and set other parameters, such as background"""
        # Add custom layout to slide
        slide = self.prs.slides.add_slide(self.title_layout)
        self.add_background(slide, self.bg_color)
        # Set title and subtitle
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        # Call to method customize_text
        customize_text(slide.shapes.title.text_frame, *self.title_font, self.title_color, bold=True, underline=True)

    def add_content_slide(self, title, content):
        """Method that add custom title to content slide and set other parameters, such as background"""
        # Add custom layout to slide
        slide = self.prs.slides.add_slide(self.content_layout)
        self.add_background(slide, self.bg_color)
        # Set title of current slide and customize it
        slide.shapes.title.text = title
        customize_text(slide.shapes.title.text_frame, *self.title_font, self.title_color, bold=True, underline=True,
                       spacing=Pt(50))
        # Set content of current slide and customize it
        slide.placeholders[1].text = "\n" + content
        customize_text(slide.placeholders[1].text_frame, *self.content_font, self.content_color, spacing=Pt(50))

    def save_presentation(self):
        """Method that call pptx save method to save presentation to the given output path"""
        self.prs.save(self.output_path)
